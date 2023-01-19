from pptx import Presentation
import pandas as pd
from pptx.dml.color import RGBColor
from pptx.util import Pt


df = pd.read_csv('Primary Academic awards.csv')
dfspecial = pd.read_csv('Primary Academic awards.csv')

df = df[df["AWARD"].str.contains("Outstanding Effort Award|Sportsman of the Year Award|Arts Award|Dux of Primary|\
                                  |Loving Hearts Strong Minds Award|Service Award|Top Academic Achiever Award|\
                                  |Sportswoman of the Year Award|All Rounder Award") == False]

dfspecial = dfspecial[dfspecial["AWARD"].str.contains(
                                "Outstanding Effort Award|Sportsman of the Year Award|Arts Award|\
                                |Dux of Primary|Loving Hearts Strong Minds Award|Service Award|\
                                |Top Academic Achiever Award|Sportswoman of the Year Award|All Rounder Award") == True]

# Combines Rows where Student number/name matches with multiple awards. Combines Award column with new line in between
dfformated = df.groupby(['STUDENT_NUMBER', 'SURNAME', 'FIRSTNAME', 'FORM',
                         'ROLL_CLASS'])['AWARD'].agg('\n'.join).reset_index()

# Sorts by surname so that slides will be in alphabetical order
dfsortsurname = dfformated.sort_values(by=['SURNAME'], ascending=True)
dfsortsurnameindex = dfsortsurname.reset_index()

# Creates custom dictionary for sorting based on Event runsheet
sort_dict = {'Hopeful Heart Award': 0, 'Knowledgeable Learner Award': 1, 'Purposeful Participant Award': 2,
             'Grace at Groves Award': 3, 'Intentional Integrity Award': 4, 'Truth and Respect Award': 5,
             'Joyful Spirit Award': 6, 'Wisdom and Discernment Award': 7, 'Perseverance Award': 8,
             'Academic Achievement Award': 9, 'Academic Achievement Award\nHopeful Heart Award': 9,
             'Academic Achievement Award\nKnowledgeable Learner Award': 9,
             'Academic Achievement Award\nPurposeful Participant Award': 9,
             'Academic Achievement Award\nGrace at Groves Award': 9,
             'Academic Achievement Award\nIntentional Integrity Award': 9,
             'Academic Achievement Award\nTruth and Respect Award': 9,
             'Academic Achievement Award\nJoyful Spirit Award': 9,
             'Academic Achievement Award\nWisdom and Discernment Award': 9,
             'Academic Achievement Award\nPerseverance Award': 9,
             'Academic Excellence Award': 10, 'Academic Excellence Award\nHopeful Heart Award': 10,
             'Academic Excellence Award\nKnowledgeable Learner Award': 10,
             'Academic Excellence Award\nPurposeful Participant Award': 10,
             'Academic Excellence Award\nGrace at Groves Award': 10,
             'Academic Excellence Award\nIntentional Integrity Award': 10,
             'Academic Excellence Award\nTruth and Respect Award': 10,
             'Academic Excellence Award\nJoyful Spirit Award': 10,
             'Academic Excellence Award\nWisdom and Discernment Award': 10,
             'Academic Excellence Award\nPerseverance Award': 10}

sort_dict_special = {'Outstanding Effort Award': 0, 'Service Award': 1, 'Top Academic Achiever Award': 2,
                     'All Rounder Award': 3, 'Loving Hearts Strong Minds Award': 4, 'Sportsman of the Year Award': 5,
                     'Sportswoman of the Year Award': 6, 'Arts Award': 7, 'Dux of Primary': 8}

# Sorts according to custom dictionary
dfsortsurnameindex.sort_values(by=['AWARD', 'SURNAME'], key=lambda x: x.map(sort_dict), inplace=True)
dfspecial.sort_values(by=['AWARD', 'FORM', 'SURNAME'], key=lambda x: x.map(sort_dict_special), inplace=True)

headerawardlist = list(dfsortsurnameindex.values.tolist())    # Create list from df so that it can run through FOR loop
with open('log.txt', 'w') as f:
    for line in headerawardlist:
        f.write(f"{line}\n")

specialawardlist = list(dfspecial.values.tolist())
with open('logspecial.txt', 'w') as f:
    for line in specialawardlist:
        f.write(f"{line}\n")

photos2022 = 'C:/Users/bwood8557/Christian Community Ministries/Groves Digital Services - Students/'
photos2021 = 'C:/Users/bwood8557/Christian Community Ministries/Groves Digital Services - Students (1)/'

prsprepC = Presentation("Awards Slide Template.pptx")
prsprepD = Presentation("Awards Slide Template.pptx")
prsprepL = Presentation("Awards Slide Template.pptx")
prsprepT = Presentation("Awards Slide Template.pptx")
prsprepR = Presentation("Awards Slide Template.pptx")
prsyear1C = Presentation("Awards Slide Template.pptx")
prsyear1W = Presentation("Awards Slide Template.pptx")
prsyear1O = Presentation("Awards Slide Template.pptx")
prsyear1F = Presentation("Awards Slide Template.pptx")
prsyear2J = Presentation("Awards Slide Template.pptx")
prsyear2F = Presentation("Awards Slide Template.pptx")
prsyear2L = Presentation("Awards Slide Template.pptx")
prsyear2B = Presentation("Awards Slide Template.pptx")
prsyear2M = Presentation("Awards Slide Template.pptx")
prsyear3M = Presentation("Awards Slide Template.pptx")
prsyear3J = Presentation("Awards Slide Template.pptx")
prsyear3T = Presentation("Awards Slide Template.pptx")
prsyear3R = Presentation("Awards Slide Template.pptx")
prsyear3F = Presentation("Awards Slide Template.pptx")
prsyear4B = Presentation("Awards Slide Template.pptx")
prsyear4E = Presentation("Awards Slide Template.pptx")
prsyear4H = Presentation("Awards Slide Template.pptx")
prsyear4K = Presentation("Awards Slide Template.pptx")
prsyear4W = Presentation("Awards Slide Template.pptx")
prsyear5R = Presentation("Awards Slide Template.pptx")
prsyear5M = Presentation("Awards Slide Template.pptx")
prsyear5C = Presentation("Awards Slide Template.pptx")
prsyear5G = Presentation("Awards Slide Template.pptx")
prsyear6E = Presentation("Awards Slide Template.pptx")
prsyear6Y = Presentation("Awards Slide Template.pptx")
prsyear6Z = Presentation("Awards Slide Template.pptx")
prsyear6P = Presentation("Awards Slide Template.pptx")
prsyear12special = Presentation("Awards Slide Template.pptx")
prsyear34special = Presentation("Awards Slide Template.pptx")
prsyear56special = Presentation("Awards Slide Template.pptx")

'''

########## Normal Awards ##########

'''


def headerawards():
    for count, row in enumerate(headerawardlist):
        idnumber = str(row[1])
        lastname = str(row[2])
        firstname = str(row[3])
        yearlevel = str(row[4])
    #    awardcategory = str(row[4])
        comments = str(row[6])
        rollclass = str(row[5])
        homeclass = 'String'

        if ' DE ' in homeclass:  # Removes DE Students from Slides
            with open('DE Students.txt', 'a') as file:
                file.write(firstname + ' ' + lastname + '\n')

        elif 'PrepC' in rollclass:
            slideprep = prsprepC.slides.add_slide(prsprepC.slide_layouts[0])
            nameplaceholder = slideprep.placeholders[11]
            imageplaceholder = slideprep.placeholders[12]
            infoplaceholder = slideprep.placeholders[13]
            nameplaceholder.text = firstname + '\n' + lastname
            infoplaceholder.text = comments
            infoplaceholder.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            try:
                infoplaceholder.text_frame.paragraphs[1].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            except Exception:
                print('')
            try:
                image = imageplaceholder.insert_picture(photos2022 + idnumber + '.jpg')
            except Exception:
                try:
                    image = imageplaceholder.insert_picture(photos2021 + idnumber + '.jpg')
                except Exception:
                    with open('nophotolog.txt', 'a') as f:
                        f.write(f'{idnumber} {firstname} {lastname} {rollclass}\n')

        elif 'PrepD' in rollclass:
            slideprep = prsprepD.slides.add_slide(prsprepD.slide_layouts[0])
            nameplaceholder = slideprep.placeholders[11]
            imageplaceholder = slideprep.placeholders[12]
            infoplaceholder = slideprep.placeholders[13]
            nameplaceholder.text = firstname + '\n' + lastname
            infoplaceholder.text = comments
            infoplaceholder.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            try:
                infoplaceholder.text_frame.paragraphs[1].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            except Exception:
                print('')
            try:
                image = imageplaceholder.insert_picture(photos2022 + idnumber + '.jpg')
            except Exception:
                try:
                    image = imageplaceholder.insert_picture(photos2021 + idnumber + '.jpg')
                except Exception:
                    with open('nophotolog.txt', 'a') as f:
                        f.write(f'{idnumber} {firstname} {lastname} {rollclass}\n')

        elif 'PrepL' in rollclass:
            slideprep = prsprepL.slides.add_slide(prsprepL.slide_layouts[0])
            nameplaceholder = slideprep.placeholders[11]
            imageplaceholder = slideprep.placeholders[12]
            infoplaceholder = slideprep.placeholders[13]
            nameplaceholder.text = firstname + '\n' + lastname
            infoplaceholder.text = comments
            infoplaceholder.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            try:
                infoplaceholder.text_frame.paragraphs[1].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            except Exception:
                print('')
            try:
                image = imageplaceholder.insert_picture(photos2022 + idnumber + '.jpg')
            except Exception:
                try:
                    image = imageplaceholder.insert_picture(photos2021 + idnumber + '.jpg')
                except Exception:
                    with open('nophotolog.txt', 'a') as f:
                        f.write(f'{idnumber} {firstname} {lastname} {rollclass}\n')

        elif 'PrepT' in rollclass:
            slideprep = prsprepT.slides.add_slide(prsprepT.slide_layouts[0])
            nameplaceholder = slideprep.placeholders[11]
            imageplaceholder = slideprep.placeholders[12]
            infoplaceholder = slideprep.placeholders[13]
            nameplaceholder.text = firstname + '\n' + lastname
            infoplaceholder.text = comments
            infoplaceholder.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            try:
                infoplaceholder.text_frame.paragraphs[1].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            except Exception:
                print('')
            try:
                image = imageplaceholder.insert_picture(photos2022 + idnumber + '.jpg')
            except Exception:
                try:
                    image = imageplaceholder.insert_picture(photos2021 + idnumber + '.jpg')
                except Exception:
                    with open('nophotolog.txt', 'a') as f:
                        f.write(f'{idnumber} {firstname} {lastname} {rollclass}\n')

        elif 'PrepR' in rollclass:
            slideprep = prsprepR.slides.add_slide(prsprepR.slide_layouts[0])
            nameplaceholder = slideprep.placeholders[11]
            imageplaceholder = slideprep.placeholders[12]
            infoplaceholder = slideprep.placeholders[13]
            nameplaceholder.text = firstname + '\n' + lastname
            infoplaceholder.text = comments
            infoplaceholder.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            try:
                infoplaceholder.text_frame.paragraphs[1].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            except Exception:
                print('')
            try:
                image = imageplaceholder.insert_picture(photos2022 + idnumber + '.jpg')
            except Exception:
                try:
                    image = imageplaceholder.insert_picture(photos2021 + idnumber + '.jpg')
                except Exception:
                    with open('nophotolog.txt', 'a') as f:
                        f.write(f'{idnumber} {firstname} {lastname} {rollclass}\n')

        elif '1C' in rollclass:
            slideyear1 = prsyear1C.slides.add_slide(prsyear1C.slide_layouts[0])
            nameplaceholder = slideyear1.placeholders[11]
            imageplaceholder = slideyear1.placeholders[12]
            infoplaceholder = slideyear1.placeholders[13]
            nameplaceholder.text = firstname + '\n' + lastname
            infoplaceholder.text = comments
            infoplaceholder.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            try:
                infoplaceholder.text_frame.paragraphs[1].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            except Exception:
                print('')
            try:
                image = imageplaceholder.insert_picture(photos2022 + idnumber + '.jpg')
            except Exception:
                try:
                    image = imageplaceholder.insert_picture(photos2021 + idnumber + '.jpg')
                except Exception:
                    with open('nophotolog.txt', 'a') as f:
                        f.write(f'{idnumber} {firstname} {lastname} {rollclass}\n')

        elif '1W' in rollclass:
            slideyear1 = prsyear1W.slides.add_slide(prsyear1W.slide_layouts[0])
            nameplaceholder = slideyear1.placeholders[11]
            imageplaceholder = slideyear1.placeholders[12]
            infoplaceholder = slideyear1.placeholders[13]
            nameplaceholder.text = firstname + '\n' + lastname
            infoplaceholder.text = comments
            infoplaceholder.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            try:
                infoplaceholder.text_frame.paragraphs[1].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            except Exception:
                print('')
            try:
                image = imageplaceholder.insert_picture(photos2022 + idnumber + '.jpg')
            except Exception:
                try:
                    image = imageplaceholder.insert_picture(photos2021 + idnumber + '.jpg')
                except Exception:
                    with open('nophotolog.txt', 'a') as f:
                        f.write(f'{idnumber} {firstname} {lastname} {rollclass}\n')

        elif '1O' in rollclass:
            slideyear1 = prsyear1O.slides.add_slide(prsyear1O.slide_layouts[0])
            nameplaceholder = slideyear1.placeholders[11]
            imageplaceholder = slideyear1.placeholders[12]
            infoplaceholder = slideyear1.placeholders[13]
            nameplaceholder.text = firstname + '\n' + lastname
            infoplaceholder.text = comments
            infoplaceholder.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            try:
                infoplaceholder.text_frame.paragraphs[1].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            except Exception:
                print('')
            try:
                image = imageplaceholder.insert_picture(photos2022 + idnumber + '.jpg')
            except Exception:
                try:
                    image = imageplaceholder.insert_picture(photos2021 + idnumber + '.jpg')
                except Exception:
                    with open('nophotolog.txt', 'a') as f:
                        f.write(f'{idnumber} {firstname} {lastname} {rollclass}\n')

        elif'1F' in rollclass:
            slideyear1 = prsyear1F.slides.add_slide(prsyear1F.slide_layouts[0])
            nameplaceholder = slideyear1.placeholders[11]
            imageplaceholder = slideyear1.placeholders[12]
            infoplaceholder = slideyear1.placeholders[13]
            nameplaceholder.text = firstname + '\n' + lastname
            infoplaceholder.text = comments
            infoplaceholder.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            try:
                infoplaceholder.text_frame.paragraphs[1].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            except Exception:
                print('')
            try:
                image = imageplaceholder.insert_picture(photos2022 + idnumber + '.jpg')
            except Exception:
                try:
                    image = imageplaceholder.insert_picture(photos2021 + idnumber + '.jpg')
                except Exception:
                    with open('nophotolog.txt', 'a') as f:
                        f.write(f'{idnumber} {firstname} {lastname} {rollclass}\n')

        elif'2J' in rollclass:
            slideyear2 = prsyear2J.slides.add_slide(prsyear2J.slide_layouts[0])
            nameplaceholder = slideyear2.placeholders[11]
            imageplaceholder = slideyear2.placeholders[12]
            infoplaceholder = slideyear2.placeholders[13]
            nameplaceholder.text = firstname + '\n' + lastname
            infoplaceholder.text = comments
            infoplaceholder.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            try:
                infoplaceholder.text_frame.paragraphs[1].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            except Exception:
                print('')
            try:
                image = imageplaceholder.insert_picture(photos2022 + idnumber + '.jpg')
            except Exception:
                try:
                    image = imageplaceholder.insert_picture(photos2021 + idnumber + '.jpg')
                except Exception:
                    with open('nophotolog.txt', 'a') as f:
                        f.write(f'{idnumber} {firstname} {lastname} {rollclass}\n')

        elif '2F' in rollclass:
            slideyear2 = prsyear2F.slides.add_slide(prsyear2F.slide_layouts[0])
            nameplaceholder = slideyear2.placeholders[11]
            imageplaceholder = slideyear2.placeholders[12]
            infoplaceholder = slideyear2.placeholders[13]
            nameplaceholder.text = firstname + '\n' + lastname
            infoplaceholder.text = comments
            infoplaceholder.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            try:
                infoplaceholder.text_frame.paragraphs[1].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            except Exception:
                print('')
            try:
                image = imageplaceholder.insert_picture(photos2022 + idnumber + '.jpg')
            except Exception:
                try:
                    image = imageplaceholder.insert_picture(photos2021 + idnumber + '.jpg')
                except Exception:
                    with open('nophotolog.txt', 'a') as f:
                        f.write(f'{idnumber} {firstname} {lastname} {rollclass}\n')

        elif '2L' in rollclass:
            slideyear2 = prsyear2L.slides.add_slide(prsyear2L.slide_layouts[0])
            nameplaceholder = slideyear2.placeholders[11]
            imageplaceholder = slideyear2.placeholders[12]
            infoplaceholder = slideyear2.placeholders[13]
            nameplaceholder.text = firstname + '\n' + lastname
            infoplaceholder.text = comments
            infoplaceholder.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            try:
                infoplaceholder.text_frame.paragraphs[1].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            except Exception:
                print('')
            try:
                image = imageplaceholder.insert_picture(photos2022 + idnumber + '.jpg')
            except Exception:
                try:
                    image = imageplaceholder.insert_picture(photos2021 + idnumber + '.jpg')
                except Exception:
                    with open('nophotolog.txt', 'a') as f:
                        f.write(f'{idnumber} {firstname} {lastname} {rollclass}\n')

        elif '2B' in rollclass:
            slideyear2 = prsyear2B.slides.add_slide(prsyear2B.slide_layouts[0])
            nameplaceholder = slideyear2.placeholders[11]
            imageplaceholder = slideyear2.placeholders[12]
            infoplaceholder = slideyear2.placeholders[13]
            nameplaceholder.text = firstname + '\n' + lastname
            infoplaceholder.text = comments
            infoplaceholder.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            try:
                infoplaceholder.text_frame.paragraphs[1].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            except Exception:
                print('')
            try:
                image = imageplaceholder.insert_picture(photos2022 + idnumber + '.jpg')
            except Exception:
                try:
                    image = imageplaceholder.insert_picture(photos2021 + idnumber + '.jpg')
                except Exception:
                    with open('nophotolog.txt', 'a') as f:
                        f.write(f'{idnumber} {firstname} {lastname} {rollclass}\n')

        elif '2M' in rollclass:
            slideyear2 = prsyear2M.slides.add_slide(prsyear2M.slide_layouts[0])
            nameplaceholder = slideyear2.placeholders[11]
            imageplaceholder = slideyear2.placeholders[12]
            infoplaceholder = slideyear2.placeholders[13]
            nameplaceholder.text = firstname + '\n' + lastname
            infoplaceholder.text = comments
            infoplaceholder.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            try:
                infoplaceholder.text_frame.paragraphs[1].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            except Exception:
                print('')
            try:
                image = imageplaceholder.insert_picture(photos2022 + idnumber + '.jpg')
            except Exception:
                try:
                    image = imageplaceholder.insert_picture(photos2021 + idnumber + '.jpg')
                except Exception:
                    with open('nophotolog.txt', 'a') as f:
                        f.write(f'{idnumber} {firstname} {lastname} {rollclass}\n')

        elif '3M' in rollclass:
            slideyear3 = prsyear3M.slides.add_slide(prsyear3M.slide_layouts[0])
            nameplaceholder = slideyear3.placeholders[11]
            imageplaceholder = slideyear3.placeholders[12]
            infoplaceholder = slideyear3.placeholders[13]
            nameplaceholder.text = firstname + '\n' + lastname
            infoplaceholder.text = comments
            infoplaceholder.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            try:
                infoplaceholder.text_frame.paragraphs[1].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            except Exception:
                print('')
            try:
                image = imageplaceholder.insert_picture(photos2022 + idnumber + '.jpg')
            except Exception:
                try:
                    image = imageplaceholder.insert_picture(photos2021 + idnumber + '.jpg')
                except Exception:
                    with open('nophotolog.txt', 'a') as f:
                        f.write(f'{idnumber} {firstname} {lastname} {rollclass}\n')

        elif '3J' in rollclass:
            slideyear3 = prsyear3J.slides.add_slide(prsyear3J.slide_layouts[0])
            nameplaceholder = slideyear3.placeholders[11]
            imageplaceholder = slideyear3.placeholders[12]
            infoplaceholder = slideyear3.placeholders[13]
            nameplaceholder.text = firstname + '\n' + lastname
            infoplaceholder.text = comments
            infoplaceholder.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            try:
                infoplaceholder.text_frame.paragraphs[1].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            except Exception:
                print('')
            try:
                image = imageplaceholder.insert_picture(photos2022 + idnumber + '.jpg')
            except Exception:
                try:
                    image = imageplaceholder.insert_picture(photos2021 + idnumber + '.jpg')
                except Exception:
                    with open('nophotolog.txt', 'a') as f:
                        f.write(f'{idnumber} {firstname} {lastname} {rollclass}\n')

        elif '3T' in rollclass:
            slideyear3 = prsyear3T.slides.add_slide(prsyear3T.slide_layouts[0])
            nameplaceholder = slideyear3.placeholders[11]
            imageplaceholder = slideyear3.placeholders[12]
            infoplaceholder = slideyear3.placeholders[13]
            nameplaceholder.text = firstname + '\n' + lastname
            infoplaceholder.text = comments
            infoplaceholder.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            try:
                infoplaceholder.text_frame.paragraphs[1].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            except Exception:
                print('')
            try:
                image = imageplaceholder.insert_picture(photos2022 + idnumber + '.jpg')
            except Exception:
                try:
                    image = imageplaceholder.insert_picture(photos2021 + idnumber + '.jpg')
                except Exception:
                    with open('nophotolog.txt', 'a') as f:
                        f.write(f'{idnumber} {firstname} {lastname} {rollclass}\n')

        elif '3R' in rollclass:
            slideyear3 = prsyear3R.slides.add_slide(prsyear3R.slide_layouts[0])
            nameplaceholder = slideyear3.placeholders[11]
            imageplaceholder = slideyear3.placeholders[12]
            infoplaceholder = slideyear3.placeholders[13]
            nameplaceholder.text = firstname + '\n' + lastname
            infoplaceholder.text = comments
            infoplaceholder.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            try:
                infoplaceholder.text_frame.paragraphs[1].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            except Exception:
                print('')
            try:
                image = imageplaceholder.insert_picture(photos2022 + idnumber + '.jpg')
            except Exception:
                try:
                    image = imageplaceholder.insert_picture(photos2021 + idnumber + '.jpg')
                except Exception:
                    with open('nophotolog.txt', 'a') as f:
                        f.write(f'{idnumber} {firstname} {lastname} {rollclass}\n')

        elif '3F' in rollclass:
            slideyear3 = prsyear3F.slides.add_slide(prsyear3F.slide_layouts[0])
            nameplaceholder = slideyear3.placeholders[11]
            imageplaceholder = slideyear3.placeholders[12]
            infoplaceholder = slideyear3.placeholders[13]
            nameplaceholder.text = firstname + '\n' + lastname
            infoplaceholder.text = comments
            infoplaceholder.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            try:
                infoplaceholder.text_frame.paragraphs[1].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            except Exception:
                print('')
            try:
                image = imageplaceholder.insert_picture(photos2022 + idnumber + '.jpg')
            except Exception:
                try:
                    image = imageplaceholder.insert_picture(photos2021 + idnumber + '.jpg')
                except Exception:
                    with open('nophotolog.txt', 'a') as f:
                        f.write(f'{idnumber} {firstname} {lastname} {rollclass}\n')

        elif '4B' in rollclass:
            slideyear4 = prsyear4B.slides.add_slide(prsyear4B.slide_layouts[0])
            nameplaceholder = slideyear4.placeholders[11]
            imageplaceholder = slideyear4.placeholders[12]
            infoplaceholder = slideyear4.placeholders[13]
            nameplaceholder.text = firstname + '\n' + lastname
            infoplaceholder.text = comments
            infoplaceholder.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            try:
                infoplaceholder.text_frame.paragraphs[1].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            except Exception:
                print('')
            try:
                image = imageplaceholder.insert_picture(photos2022 + idnumber + '.jpg')
            except Exception:
                try:
                    image = imageplaceholder.insert_picture(photos2021 + idnumber + '.jpg')
                except Exception:
                    with open('nophotolog.txt', 'a') as f:
                        f.write(f'{idnumber} {firstname} {lastname} {rollclass}\n')

        elif '4E' in rollclass:
            slideyear4 = prsyear4E.slides.add_slide(prsyear4E.slide_layouts[0])
            nameplaceholder = slideyear4.placeholders[11]
            imageplaceholder = slideyear4.placeholders[12]
            infoplaceholder = slideyear4.placeholders[13]
            nameplaceholder.text = firstname + '\n' + lastname
            infoplaceholder.text = comments
            infoplaceholder.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            try:
                infoplaceholder.text_frame.paragraphs[1].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            except Exception:
                print('')
            try:
                image = imageplaceholder.insert_picture(photos2022 + idnumber + '.jpg')
            except Exception:
                try:
                    image = imageplaceholder.insert_picture(photos2021 + idnumber + '.jpg')
                except Exception:
                    with open('nophotolog.txt', 'a') as f:
                        f.write(f'{idnumber} {firstname} {lastname} {rollclass}\n')

        elif '4H' in rollclass:
            slideyear4 = prsyear4H.slides.add_slide(prsyear4H.slide_layouts[0])
            nameplaceholder = slideyear4.placeholders[11]
            imageplaceholder = slideyear4.placeholders[12]
            infoplaceholder = slideyear4.placeholders[13]
            nameplaceholder.text = firstname + '\n' + lastname
            infoplaceholder.text = comments
            infoplaceholder.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            try:
                infoplaceholder.text_frame.paragraphs[1].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            except Exception:
                print('')
            try:
                image = imageplaceholder.insert_picture(photos2022 + idnumber + '.jpg')
            except Exception:
                try:
                    image = imageplaceholder.insert_picture(photos2021 + idnumber + '.jpg')
                except Exception:
                    with open('nophotolog.txt', 'a') as f:
                        f.write(f'{idnumber} {firstname} {lastname} {rollclass}\n')

        elif '4K' in rollclass:
            slideyear4 = prsyear4K.slides.add_slide(prsyear4K.slide_layouts[0])
            nameplaceholder = slideyear4.placeholders[11]
            imageplaceholder = slideyear4.placeholders[12]
            infoplaceholder = slideyear4.placeholders[13]
            nameplaceholder.text = firstname + '\n' + lastname
            infoplaceholder.text = comments
            infoplaceholder.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            try:
                infoplaceholder.text_frame.paragraphs[1].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            except Exception:
                print('')
            try:
                image = imageplaceholder.insert_picture(photos2022 + idnumber + '.jpg')
            except Exception:
                try:
                    image = imageplaceholder.insert_picture(photos2021 + idnumber + '.jpg')
                except Exception:
                    with open('nophotolog.txt', 'a') as f:
                        f.write(f'{idnumber} {firstname} {lastname} {rollclass}\n')

        elif '4W' in rollclass:
            slideyear4 = prsyear4W.slides.add_slide(prsyear4W.slide_layouts[0])
            nameplaceholder = slideyear4.placeholders[11]
            imageplaceholder = slideyear4.placeholders[12]
            infoplaceholder = slideyear4.placeholders[13]
            nameplaceholder.text = firstname + '\n' + lastname
            infoplaceholder.text = comments
            infoplaceholder.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            try:
                infoplaceholder.text_frame.paragraphs[1].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            except Exception:
                print('')
            try:
                image = imageplaceholder.insert_picture(photos2022 + idnumber + '.jpg')
            except Exception:
                try:
                    image = imageplaceholder.insert_picture(photos2021 + idnumber + '.jpg')
                except Exception:
                    with open('nophotolog.txt', 'a') as f:
                        f.write(f'{idnumber} {firstname} {lastname} {rollclass}\n')

        elif '5R' in rollclass:
            slideyear5 = prsyear5R.slides.add_slide(prsyear5R.slide_layouts[0])
            nameplaceholder = slideyear5.placeholders[11]
            imageplaceholder = slideyear5.placeholders[12]
            infoplaceholder = slideyear5.placeholders[13]
            nameplaceholder.text = firstname + '\n' + lastname
            infoplaceholder.text = comments
            infoplaceholder.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            try:
                infoplaceholder.text_frame.paragraphs[1].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            except Exception:
                print('')
            try:
                image = imageplaceholder.insert_picture(photos2022 + idnumber + '.jpg')
            except Exception:
                try:
                    image = imageplaceholder.insert_picture(photos2021 + idnumber + '.jpg')
                except Exception:
                    with open('nophotolog.txt', 'a') as f:
                        f.write(f'{idnumber} {firstname} {lastname} {rollclass}\n')

        elif '5M' in rollclass:
            slideyear5 = prsyear5M.slides.add_slide(prsyear5M.slide_layouts[0])
            nameplaceholder = slideyear5.placeholders[11]
            imageplaceholder = slideyear5.placeholders[12]
            infoplaceholder = slideyear5.placeholders[13]
            nameplaceholder.text = firstname + '\n' + lastname
            infoplaceholder.text = comments
            infoplaceholder.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            try:
                infoplaceholder.text_frame.paragraphs[1].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            except Exception:
                print('')
            try:
                image = imageplaceholder.insert_picture(photos2022 + idnumber + '.jpg')
            except Exception:
                try:
                    image = imageplaceholder.insert_picture(photos2021 + idnumber + '.jpg')
                except Exception:
                    with open('nophotolog.txt', 'a') as f:
                        f.write(f'{idnumber} {firstname} {lastname} {rollclass}\n')

        elif '5C' in rollclass:
            slideyear5 = prsyear5C.slides.add_slide(prsyear5C.slide_layouts[0])
            nameplaceholder = slideyear5.placeholders[11]
            imageplaceholder = slideyear5.placeholders[12]
            infoplaceholder = slideyear5.placeholders[13]
            nameplaceholder.text = firstname + '\n' + lastname
            infoplaceholder.text = comments
            infoplaceholder.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            try:
                infoplaceholder.text_frame.paragraphs[1].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            except Exception:
                print('')
            try:
                image = imageplaceholder.insert_picture(photos2022 + idnumber + '.jpg')
            except Exception:
                try:
                    image = imageplaceholder.insert_picture(photos2021 + idnumber + '.jpg')
                except Exception:
                    with open('nophotolog.txt', 'a') as f:
                        f.write(f'{idnumber} {firstname} {lastname} {rollclass}\n')

        elif '5G' in rollclass:
            slideyear5 = prsyear5G.slides.add_slide(prsyear5G.slide_layouts[0])
            nameplaceholder = slideyear5.placeholders[11]
            imageplaceholder = slideyear5.placeholders[12]
            infoplaceholder = slideyear5.placeholders[13]
            nameplaceholder.text = firstname + '\n' + lastname
            infoplaceholder.text = comments
            infoplaceholder.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            try:
                infoplaceholder.text_frame.paragraphs[1].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            except Exception:
                print('')
            try:
                image = imageplaceholder.insert_picture(photos2022 + idnumber + '.jpg')
            except Exception:
                try:
                    image = imageplaceholder.insert_picture(photos2021 + idnumber + '.jpg')
                except Exception:
                    with open('nophotolog.txt', 'a') as f:
                        f.write(f'{idnumber} {firstname} {lastname} {rollclass}\n')

        elif '6E' in rollclass:
            slideyear6 = prsyear6E.slides.add_slide(prsyear6E.slide_layouts[0])
            nameplaceholder = slideyear6.placeholders[11]
            imageplaceholder = slideyear6.placeholders[12]
            infoplaceholder = slideyear6.placeholders[13]
            nameplaceholder.text = firstname + '\n' + lastname
            infoplaceholder.text = comments
            infoplaceholder.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            try:
                infoplaceholder.text_frame.paragraphs[1].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            except Exception:
                print('')
            try:
                image = imageplaceholder.insert_picture(photos2022 + idnumber + '.jpg')
            except Exception:
                try:
                    image = imageplaceholder.insert_picture(photos2021 + idnumber + '.jpg')
                except Exception:
                    with open('nophotolog.txt', 'a') as f:
                        f.write(f'{idnumber} {firstname} {lastname} {rollclass}\n')

        elif '6Y' in rollclass:
            slideyear6 = prsyear6Y.slides.add_slide(prsyear6Y.slide_layouts[0])
            nameplaceholder = slideyear6.placeholders[11]
            imageplaceholder = slideyear6.placeholders[12]
            infoplaceholder = slideyear6.placeholders[13]
            nameplaceholder.text = firstname + '\n' + lastname
            infoplaceholder.text = comments
            infoplaceholder.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            try:
                infoplaceholder.text_frame.paragraphs[1].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            except Exception:
                print('')
            try:
                image = imageplaceholder.insert_picture(photos2022 + idnumber + '.jpg')
            except Exception:
                try:
                    image = imageplaceholder.insert_picture(photos2021 + idnumber + '.jpg')
                except Exception:
                    with open('nophotolog.txt', 'a') as f:
                        f.write(f'{idnumber} {firstname} {lastname} {rollclass}\n')

        elif '6Z' in rollclass:
            slideyear6 = prsyear6Z.slides.add_slide(prsyear6Z.slide_layouts[0])
            nameplaceholder = slideyear6.placeholders[11]
            imageplaceholder = slideyear6.placeholders[12]
            infoplaceholder = slideyear6.placeholders[13]
            nameplaceholder.text = firstname + '\n' + lastname
            infoplaceholder.text = comments
            infoplaceholder.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            try:
                infoplaceholder.text_frame.paragraphs[1].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            except Exception:
                print('')
            try:
                image = imageplaceholder.insert_picture(photos2022 + idnumber + '.jpg')
            except Exception:
                try:
                    image = imageplaceholder.insert_picture(photos2021 + idnumber + '.jpg')
                except Exception:
                    with open('nophotolog.txt', 'a') as f:
                        f.write(f'{idnumber} {firstname} {lastname} {rollclass}\n')

        elif '6P' in rollclass:
            slideyear6 = prsyear6P.slides.add_slide(prsyear6P.slide_layouts[0])
            nameplaceholder = slideyear6.placeholders[11]
            imageplaceholder = slideyear6.placeholders[12]
            infoplaceholder = slideyear6.placeholders[13]
            nameplaceholder.text = firstname + '\n' + lastname
            infoplaceholder.text = comments
            infoplaceholder.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            try:
                infoplaceholder.text_frame.paragraphs[1].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            except Exception:
                print('')
            try:
                image = imageplaceholder.insert_picture(photos2022 + idnumber + '.jpg')
            except Exception:
                try:
                    image = imageplaceholder.insert_picture(photos2021 + idnumber + '.jpg')
                except Exception:
                    with open('nophotolog.txt', 'a') as f:
                        f.write(f'{idnumber} {firstname} {lastname} {rollclass}\n')


'''

########### Special Awards #############

'''


def specialawards():
    for count, row in enumerate(specialawardlist):
        idnumber = str(row[0])
        lastname = str(row[1])
        firstname = str(row[2])
        yearlevel = str(row[3])
    #    awardcategory = str(row[4])
        comments = str(row[6])
    #    rollclass = str(row[6])
        homeclass = 'String'

        if ' DE ' in homeclass:  # Removes DE Students from Slides
            with open('DE Students.txt', 'a') as file:
                file.write(firstname + ' ' + lastname + '\n')

        elif 'Preparatory' in yearlevel:
            slideprep = prsprep.slides.add_slide(prsprep.slide_layouts[1])
            infoplaceholder = slideprep.placeholders[11]
            imageplaceholder = slideprep.placeholders[12]
            nameplaceholder = slideprep.placeholders[13]
            nameplaceholder.text = firstname + '\n' + lastname
            infoplaceholder.text = comments
            nameplaceholder.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            nameplaceholder.text_frame.paragraphs[0].font.size = Pt(48)
            nameplaceholder.text_frame.paragraphs[1].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            nameplaceholder.text_frame.paragraphs[1].font.size = Pt(48)
            try:
                image = imageplaceholder.insert_picture(photos2022 + idnumber + '.jpg')
            except Exception:
                try:
                    image = imageplaceholder.insert_picture(photos2021 + idnumber + '.jpg')
                except Exception:
                    with open('nophotolog.txt', 'a') as f:
                        f.write(idnumber + ' ' + firstname + ' ' + lastname + '\n')

        elif 'Year 01' in yearlevel or 'Year 02' in yearlevel:
            slideyear12 = prsyear12special.slides.add_slide(prsyear12special.slide_layouts[1])
            infoplaceholder = slideyear12.placeholders[11]
            imageplaceholder = slideyear12.placeholders[12]
            nameplaceholder = slideyear12.placeholders[13]
            nameplaceholder.text = firstname + '\n' + lastname
            infoplaceholder.text = comments
            nameplaceholder.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            nameplaceholder.text_frame.paragraphs[0].font.size = Pt(48)
            nameplaceholder.text_frame.paragraphs[1].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            nameplaceholder.text_frame.paragraphs[1].font.size = Pt(48)
            try:
                image = imageplaceholder.insert_picture(photos2022 + idnumber + '.jpg')
            except Exception:
                try:
                    image = imageplaceholder.insert_picture(photos2021 + idnumber + '.jpg')
                except Exception:
                    with open('nophotolog.txt', 'a') as f:
                        f.write(idnumber + ' ' + firstname + ' ' + lastname + '\n')

        elif 'Year 03' in yearlevel or 'Year 04' in yearlevel:
            slideyear34 = prsyear34special.slides.add_slide(prsyear34special.slide_layouts[1])
            infoplaceholder = slideyear34.placeholders[11]
            imageplaceholder = slideyear34.placeholders[12]
            nameplaceholder = slideyear34.placeholders[13]
            nameplaceholder.text = firstname + '\n' + lastname
            infoplaceholder.text = comments
            nameplaceholder.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            nameplaceholder.text_frame.paragraphs[0].font.size = Pt(48)
            nameplaceholder.text_frame.paragraphs[1].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            nameplaceholder.text_frame.paragraphs[1].font.size = Pt(48)
            try:
                image = imageplaceholder.insert_picture(photos2022 + idnumber + '.jpg')
            except Exception:
                try:
                    image = imageplaceholder.insert_picture(photos2021 + idnumber + '.jpg')
                except Exception:
                    with open('nophotolog.txt', 'a') as f:
                        f.write(idnumber + ' ' + firstname + ' ' + lastname + '\n')

        elif 'Year 05' in yearlevel or 'Year 06' in yearlevel:
            slideyear56 = prsyear56special.slides.add_slide(prsyear56special.slide_layouts[1])
            infoplaceholder = slideyear56.placeholders[11]
            imageplaceholder = slideyear56.placeholders[12]
            nameplaceholder = slideyear56.placeholders[13]
            nameplaceholder.text = firstname + '\n' + lastname
            infoplaceholder.text = comments
            nameplaceholder.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            nameplaceholder.text_frame.paragraphs[0].font.size = Pt(48)
            nameplaceholder.text_frame.paragraphs[1].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            nameplaceholder.text_frame.paragraphs[1].font.size = Pt(48)
            try:
                image = imageplaceholder.insert_picture(photos2022 + idnumber + '.jpg')
            except Exception:
                try:
                    image = imageplaceholder.insert_picture(photos2021 + idnumber + '.jpg')
                except Exception:
                    with open('nophotolog.txt', 'a') as f:
                        f.write(idnumber + ' ' + firstname + ' ' + lastname + '\n')


headerawards()
specialawards()


prsprepC.save('Prep C Awards.pptx')
prsprepD.save('Prep D Awards.pptx')
prsprepL.save('Prep L Awards.pptx')
prsprepR.save('Prep R Awards.pptx')
prsprepT.save('Prep T Awards.pptx')
prsyear1C.save('Year1C Awards.pptx')
prsyear1W.save('Year1W Awards.pptx')
prsyear1O.save('Year1O Awards.pptx')
prsyear1F.save('Year1F Awards.pptx')
prsyear2J.save('Year2J Awards.pptx')
prsyear2F.save('Year2F Awards.pptx')
prsyear2L.save('Year2L Awards.pptx')
prsyear2B.save('Year2B Awards.pptx')
prsyear2M.save('Year2M Awards.pptx')
prsyear3M.save('Year3M Awards.pptx')
prsyear3J.save('Year3J Awards.pptx')
prsyear3T.save('Year3T Awards.pptx')
prsyear3R.save('Year3R Awards.pptx')
prsyear3F.save('Year3F Awards.pptx')
prsyear4B.save('Year4B Awards.pptx')
prsyear4E.save('Year4E Awards.pptx')
prsyear4H.save('Year4H Awards.pptx')
prsyear4K.save('Year4K Awards.pptx')
prsyear4W.save('Year4W Awards.pptx')
prsyear5R.save('Year5R Awards.pptx')
prsyear5M.save('Year5M Awards.pptx')
prsyear5C.save('Year5C Awards.pptx')
prsyear5G.save('Year5G Awards.pptx')
prsyear6E.save('Year6E Awards.pptx')
prsyear6Y.save('Year6Y Awards.pptx')
prsyear6Z.save('Year6Z Awards.pptx')
prsyear6P.save('Year6P Awards.pptx')
prsyear12special.save('Year12 Special Awards.pptx')
prsyear34special.save('Year34 Special Awards.pptx')
prsyear56special.save('Year56 Special Awards.pptx')
